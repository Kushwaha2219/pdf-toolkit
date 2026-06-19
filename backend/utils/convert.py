"""Convert between PDF and other formats (Word, PowerPoint, Excel, images)."""

import os
import shutil
import subprocess
import zipfile

import fitz  # PyMuPDF
import img2pdf
from pdf2docx import Converter


# --------------------------------------------------------------------------- #
# PDF -> Word
# --------------------------------------------------------------------------- #
def pdf_to_docx(input_path, output_path):
    """Convert a PDF into an editable Word (.docx) document."""
    cv = Converter(input_path)
    try:
        cv.convert(output_path, start=0, end=None)
    finally:
        cv.close()
    return output_path


# --------------------------------------------------------------------------- #
# Images <-> PDF
# --------------------------------------------------------------------------- #
def images_to_pdf(image_paths, output_path):
    """Combine one or more images (JPG/PNG/...) into a single PDF."""
    if not image_paths:
        raise ValueError("At least one image is required.")
    with open(output_path, "wb") as f:
        f.write(img2pdf.convert(image_paths))
    return output_path


def pdf_to_images(input_path, output_dir, fmt="png", dpi=150):
    """Render each PDF page to an image file. Returns the list of image paths."""
    os.makedirs(output_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(input_path))[0]
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    doc = fitz.open(input_path)
    paths = []
    try:
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=matrix)
            out = os.path.join(output_dir, f"{base}_page{i}.{fmt}")
            pix.save(out)
            paths.append(out)
    finally:
        doc.close()
    return paths


def pdf_to_jpg_zip(input_path, output_dir, output_zip, dpi=150):
    """Render each PDF page to a JPG and bundle them into a ZIP."""
    images = pdf_to_images(input_path, output_dir, fmt="jpg", dpi=dpi)
    with zipfile.ZipFile(output_zip, "w", zipfile.ZIP_DEFLATED) as zf:
        for path in images:
            zf.write(path, arcname=os.path.basename(path))
    return output_zip


# --------------------------------------------------------------------------- #
# PDF -> PowerPoint  (each page becomes a full-bleed slide image)
# --------------------------------------------------------------------------- #
def pdf_to_pptx(input_path, output_path, dpi=150):
    """Convert a PDF into a PPTX where each page is rendered onto one slide."""
    from pptx import Presentation
    from pptx.util import Emu

    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    tmp_dir = output_path + "_imgs"
    os.makedirs(tmp_dir, exist_ok=True)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # fully blank
    EMU_PER_PT = 12700  # 1 point = 12700 EMU

    doc = fitz.open(input_path)
    try:
        for i, page in enumerate(doc, start=1):
            # Size the slide to the page (points -> EMU).
            prs.slide_width = Emu(int(page.rect.width * EMU_PER_PT))
            prs.slide_height = Emu(int(page.rect.height * EMU_PER_PT))

            pix = page.get_pixmap(matrix=matrix)
            img_path = os.path.join(tmp_dir, f"page{i}.png")
            pix.save(img_path)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path, 0, 0, width=prs.slide_width, height=prs.slide_height
            )
    finally:
        doc.close()

    prs.save(output_path)
    shutil.rmtree(tmp_dir, ignore_errors=True)
    return output_path


# --------------------------------------------------------------------------- #
# PDF -> Excel  (extract tables; fall back to page text)
# --------------------------------------------------------------------------- #
def pdf_to_xlsx(input_path, output_path):
    """Extract tables from a PDF into an .xlsx workbook (one sheet per page).

    If a page has no detectable table, its text lines are written instead so
    nothing is silently lost.
    """
    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)  # start with no sheets

    with pdfplumber.open(input_path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            ws = wb.create_sheet(title=f"Page {page_no}"[:31])
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table:
                        ws.append([("" if c is None else c) for c in row])
                    ws.append([])  # blank row between tables
            else:
                text = page.extract_text() or ""
                for line in text.splitlines():
                    ws.append([line])

    if not wb.sheetnames:
        wb.create_sheet(title="Empty")

    wb.save(output_path)
    return output_path


# --------------------------------------------------------------------------- #
# Office (Word / PowerPoint / Excel) -> PDF  via LibreOffice headless
# --------------------------------------------------------------------------- #
def _find_soffice():
    """Locate the LibreOffice binary across Linux/Windows installs."""
    for name in ("soffice", "libreoffice"):
        path = shutil.which(name)
        if path:
            return path
    # Common Windows install location.
    win = r"C:\Program Files\LibreOffice\program\soffice.exe"
    if os.path.exists(win):
        return win
    return None


def office_to_pdf(input_path, output_dir):
    """Convert a Word/PPT/Excel file to PDF using LibreOffice headless.

    Returns the path to the produced PDF. Raises ValueError with a friendly
    message if LibreOffice isn't available (e.g. local Windows without it).
    """
    soffice = _find_soffice()
    if not soffice:
        raise ValueError(
            "LibreOffice is not installed on this server, so Office-to-PDF "
            "conversion is unavailable here. (It is installed in the deployed "
            "Docker image.)"
        )

    os.makedirs(output_dir, exist_ok=True)
    result = subprocess.run(
        [
            soffice, "--headless", "--norestore", "--convert-to", "pdf",
            "--outdir", output_dir, input_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise ValueError(f"Conversion failed: {result.stderr.strip() or 'unknown error'}")

    base = os.path.splitext(os.path.basename(input_path))[0]
    produced = os.path.join(output_dir, base + ".pdf")
    if not os.path.exists(produced):
        raise ValueError("Conversion produced no output file.")
    return produced
